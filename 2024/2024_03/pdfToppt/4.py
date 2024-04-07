from pptx import Presentation

# 打开PPT文件
ppt = Presentation('2222.ppt')

# 遍历每个幻灯片
for slide in ppt.slides:
    # 遍历每个形状（包括图片等）在幻灯片中
    for shape in slide.shapes:
        # 筛选出图片类型的形状
        if shape.shape_type == 13:  # 假设图片类型对应的编号是13
            # 获得图片的大小和位置信息
            top = shape.top
            left = shape.left
            width = shape.width
            height = shape.height
            
            # 在这里可以对位置和大小进行处理，应用到其他PPT等
            
            # 输出信息（可根据需要进行操作）
            print("Top:", top)
            print("Left:", left)
            print("Width:", width)
            print("Height:", height)