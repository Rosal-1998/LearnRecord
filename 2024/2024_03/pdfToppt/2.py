import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io

def pdf_page_to_ppt_sections_high_dpi(pdf_path, ppt_path, dpi=1000):
    doc = fitz.open(pdf_path)
    prs = Presentation()

    # 假设的每个部分的大小和位置（这里需要根据实际情况调整）
    sections = [
        [(30,5), (200, 30)], # 第1部分的位置
        [(7,35), (223, 159)], # 第2部分
        # 根据实际布局继续添加其他部分的位置
    ]

    for i, page in enumerate(doc):
        slide_layout = prs.slide_layouts[5] # 使用第五个幻灯片布局（空白布局）
        slide = prs.slides.add_slide(slide_layout)

        for section in sections:
            # 截取指定区域
            clip_rect = fitz.Rect(*section[0], *section[1])
            pix = page.get_pixmap(clip=clip_rect, dpi=dpi)
            img_data = io.BytesIO(pix.tobytes("png"))

            # 添加图片至幻灯片
            left = Inches(0)
            top = Inches(0)
            width = Inches(pix.width / dpi) # 计算图片宽度（英寸）
            height = Inches(pix.height / dpi) # 计算图片高度（英寸）
            slide.shapes.add_picture(img_data, left, top, width, height)
            
    prs.save(ppt_path)

pdf_path = '2.pdf'  # PDF文件路径
ppt_path = 'output_ppt_file.pptx'  # 输出PPT文件路径
pdf_page_to_ppt_sections_high_dpi(pdf_path, ppt_path, dpi=1000)