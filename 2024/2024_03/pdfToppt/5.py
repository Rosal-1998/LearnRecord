from pptx import Presentation
# 调整位置
def update_ppt(ppt_file, small_top, small_left, small_width, small_height, large_top, large_left, large_width, large_height):
    ppt = Presentation(ppt_file)
    
    for slide in ppt.slides:
        # 遍历每页中的两个图片形状
        image_shapes = [(i, shape) for i, shape in enumerate(slide.shapes) if shape.shape_type == 13]
        if len(image_shapes) != 2:  # 如果不是两个图片，则跳过该页
            continue
        
        # 根据图片大小判断出较小和较大的图片形状
        if image_shapes[0][1].width <= image_shapes[1][1].width:
            small_index, small_shape = image_shapes[0]
            large_index, large_shape = image_shapes[1]
        else:
            small_index, small_shape = image_shapes[1]
            large_index, large_shape = image_shapes[0]
        
        # 更新较小图片的位置和大小
        small_shape.top = small_top
        small_shape.left = small_left
        small_shape.width = small_width
        small_shape.height = small_height
        
        # 更新较大图片的位置和大小
        large_shape.top = large_top
        large_shape.left = large_left
        large_shape.width = large_width
        large_shape.height = large_height
    
    # 保存修改后的PPT文件
    ppt.save('c99.ppt')

# 指定要处理的PPT文件路径和文件名
ppt_file = 'c99.ppt'

# 指定给定的位置和大小值
small_top = 122555
small_left = 2151380
small_width = 4840605
small_height = 713740

large_top = 1615440
large_left = 724535
large_width = 8376285
large_height = 4810125



# 调用函数进行PPT处理
update_ppt(ppt_file, small_top, small_left, small_width, small_height, large_top, large_left, large_width, large_height)