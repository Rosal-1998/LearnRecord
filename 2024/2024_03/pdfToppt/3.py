import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io
#生成ppt
def pdf_page_to_ppt_sections_high_dpi(pdf_path, ppt_path, dpi=1000):
    doc = fitz.open(pdf_path)
    prs = Presentation()

    sections = [
        [(30, 5), (200, 30)],   # 第1部分的位置和大小
        [(6,33), (205, 148)]   # 第2部分的位置和大小
        # 根据实际布局继续添加其他部分的位置和大小
    ]

    for i, page in enumerate(doc):
        slide_layout = prs.slide_layouts[6]  # 使用第6个幻灯片布局（空白布局）
        slide = prs.slides.add_slide(slide_layout)

        for j, section in enumerate(sections):
            # 截取指定区域
            clip_rect = fitz.Rect(*section[0], *section[1])
            pix = page.get_pixmap(clip=clip_rect, dpi=dpi)
            img_data = io.BytesIO(pix.tobytes("png"))

            if j == 0:  # 第一部分
                left = 0
                top = 0
                width = pix.width * 1.5
                height = pix.height * 1.5
            else:  # 第二部分
                left = width
                top = 0
                width = pix.width * 2
                height = pix.height * 2

            # 添加图片至幻灯片
            slide.shapes.add_picture(img_data, left=Inches(left/dpi), top=Inches(top/dpi), width=Inches(width/dpi), height=Inches(height/dpi))

        prs.slide_width = Inches(width/dpi)
        prs.slide_height = Inches(height/dpi)

    prs.save(ppt_path)

pdf_path = 'C9.pdf'  # PDF文件路径
ppt_path = 'C9.ppt'  # 输出PPT文件路径
pdf_page_to_ppt_sections_high_dpi(pdf_path, ppt_path, dpi=1000)